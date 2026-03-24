import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

YELLOW = RGBColor(0xFF, 0xFF, 0x00)
BLACK  = RGBColor(0x00, 0x00, 0x00)
GREY   = RGBColor(0x99, 0x99, 0x99)


def download_image(url: str, filename: str) -> str | None:
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        with open(filename, "wb") as f:
            f.write(response.content)
        return filename
    except requests.RequestException as e:
        print(f"  [warn] Could not download {url}: {e}")
        return None


def cleanup_temp_images(paths: list[str]) -> None:
    for path in paths:
        try:
            os.remove(path)
        except OSError:
            pass


def add_question_slide(prs: Presentation, q: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Black background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BLACK

    # Question label: "Q7.57."
    id_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(3), Inches(0.5))
    id_tf = id_box.text_frame
    id_tf.text = f"Q{q['question_id']}."
    id_tf.paragraphs[0].runs[0].font.size = Pt(20)
    id_tf.paragraphs[0].runs[0].font.bold = True
    id_tf.paragraphs[0].runs[0].font.color.rgb = YELLOW

    has_images = bool(q["images"])

    if has_images:
        # Two-column: question text left, images right
        text_w = Inches(5.2)
        text_h = Inches(6.3)
        img_x  = Inches(5.8)
        img_w  = Inches(3.8)
        img_h  = Inches(6.3) / max(len(q["images"]), 1) - Inches(0.1)
    else:
        text_w = Inches(9.2)
        text_h = Inches(6.3)

    # Question text
    txt_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), text_w, text_h)
    tf = txt_box.text_frame
    tf.word_wrap = True
    tf.text = q["raw"]
    for para in tf.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = YELLOW

    if not has_images:
        return

    # Download and place images
    temp_files = []
    loaded_any = False

    for j, img_url in enumerate(q["images"]):
        img_path = f"temp_img_{q['question_id']}_{j}.png"
        downloaded = download_image(img_url, img_path)

        if downloaded:
            temp_files.append(downloaded)
            y_pos = Inches(0.9) + j * (img_h + Inches(0.1))
            try:
                slide.shapes.add_picture(downloaded, img_x, y_pos, width=img_w, height=img_h)
                loaded_any = True
            except Exception as e:
                print(f"  [warn] Could not add image to slide: {e}")

    cleanup_temp_images(temp_files)

    # Fallback label if all images failed
    if not loaded_any:
        fb = slide.shapes.add_textbox(img_x, Inches(3.5), Inches(3.8), Inches(0.6))
        fb.text_frame.text = "[ Diagram unavailable ]"
        fb.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        fb.text_frame.paragraphs[0].runs[0].font.color.rgb = GREY


def create_ppt(questions: list[dict], output_file: str = "questions.pptx") -> None:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for q in questions:
        add_question_slide(prs, q)

    prs.save(output_file)
    print(f"Saved: {output_file}  ({len(questions)} slides)")

#OLD
#-------------------------------------------------------------------------
# import requests
# from pptx import Presentation
# from pptx.util import Inches


# def download_image(url, filename):
#     try:
#         response = requests.get(url)
#         if response.status_code == 200:
#             with open(filename, "wb") as f:
#                 f.write(response.content)
#             return filename
#     except:
#         pass
#     return None


# def create_ppt(questions, output_file="questions.pptx"):
#     prs = Presentation()

#     for i, q in enumerate(questions):
#         slide_layout = prs.slide_layouts[1]  # title + content
#         slide = prs.slides.add_slide(slide_layout)

#         # Title
#         slide.shapes.title.text = q["question_id"]

#         # Content
#         content = slide.placeholders[1]
#         content.text = q["raw"]

#         # Add images (if any)
#         y_offset = 3  # start lower to avoid overlap

#         for j, img_url in enumerate(q["images"]):
#             img_path = f"temp_img_{i}_{j}.png"
#             downloaded = download_image(img_url, img_path)

#             if downloaded:
#                 slide.shapes.add_picture(
#                     img_path,
#                     Inches(1),
#                     Inches(y_offset),
#                     width=Inches(5)
#                 )
#                 y_offset += 2.5  # stack images

#     prs.save(output_file)
#     print(f"PPT saved as {output_file}")